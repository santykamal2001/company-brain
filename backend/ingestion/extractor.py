"""
Multi-format text extraction. Tries native text first; falls back to OCR for image pages.
Returns plain text. File type is determined by mime type, not extension.
"""
from __future__ import annotations

import io
import mimetypes
from pathlib import Path


def extract(file_path: str | Path, mime_type: str | None = None) -> str:
    path = Path(file_path)
    mt = mime_type or mimetypes.guess_type(str(path))[0] or ""

    if mt == "application/pdf" or path.suffix.lower() == ".pdf":
        return _extract_pdf(path)
    if mt in ("application/vnd.openxmlformats-officedocument.wordprocessingml.document",) or path.suffix.lower() == ".docx":
        return _extract_docx(path)
    if mt in ("application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",) or path.suffix.lower() == ".xlsx":
        return _extract_xlsx(path)
    if mt in ("application/vnd.openxmlformats-officedocument.presentationml.presentation",) or path.suffix.lower() == ".pptx":
        return _extract_pptx(path)
    if mt == "text/csv" or path.suffix.lower() == ".csv":
        return _extract_csv(path)
    # Fallback: plain text
    return _extract_text(path)


def _extract_pdf(path: Path) -> str:
    import fitz  # PyMuPDF

    doc = fitz.open(str(path))
    parts: list[str] = []
    ocr_pages: list[int] = []

    for i, page in enumerate(doc):
        text = page.get_text("text").strip()
        if text:
            parts.append(text)
        else:
            ocr_pages.append(i)

    if ocr_pages:
        ocr_text = _ocr_pages(path, ocr_pages)
        parts.extend(ocr_text)

    doc.close()
    return "\n\n".join(parts)


def _ocr_pages(path: Path, page_indices: list[int]) -> list[str]:
    try:
        from paddleocr import PaddleOCR
        import fitz
        import numpy as np

        ocr = PaddleOCR(use_angle_cls=True, lang="en", show_log=False)
        doc = fitz.open(str(path))
        results: list[str] = []
        for i in page_indices:
            page = doc[i]
            pix = page.get_pixmap(dpi=200)
            img_array = np.frombuffer(pix.samples, dtype=np.uint8).reshape(pix.h, pix.w, pix.n)
            result = ocr.ocr(img_array, cls=True)
            if result and result[0]:
                page_text = "\n".join(line[1][0] for line in result[0] if line[1])
                results.append(page_text)
        doc.close()
        return results
    except Exception:
        return []


def _extract_docx(path: Path) -> str:
    from docx import Document

    doc = Document(str(path))
    parts: list[str] = []
    for para in doc.paragraphs:
        if para.text.strip():
            parts.append(para.text)
    for table in doc.tables:
        for row in table.rows:
            row_text = " | ".join(cell.text.strip() for cell in row.cells if cell.text.strip())
            if row_text:
                parts.append(row_text)
    return "\n".join(parts)


def _extract_xlsx(path: Path) -> str:
    import openpyxl

    wb = openpyxl.load_workbook(str(path), data_only=True, read_only=True)
    parts: list[str] = []
    for sheet in wb.worksheets:
        parts.append(f"[Sheet: {sheet.title}]")
        for row in sheet.iter_rows(values_only=True):
            row_text = " | ".join(str(v) for v in row if v is not None)
            if row_text:
                parts.append(row_text)
    wb.close()
    return "\n".join(parts)


def _extract_pptx(path: Path) -> str:
    from pptx import Presentation

    prs = Presentation(str(path))
    parts: list[str] = []
    for i, slide in enumerate(prs.slides, 1):
        slide_texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    t = para.text.strip()
                    if t:
                        slide_texts.append(t)
        if slide_texts:
            parts.append(f"[Slide {i}]\n" + "\n".join(slide_texts))
    return "\n\n".join(parts)


def _extract_csv(path: Path) -> str:
    import csv

    for encoding in ("utf-8", "latin-1", "cp1252"):
        try:
            with path.open("r", encoding=encoding, newline="") as f:
                reader = csv.reader(f)
                rows = [" | ".join(row) for row in reader if any(cell.strip() for cell in row)]
            return "\n".join(rows)
        except (UnicodeDecodeError, Exception):
            continue
    return ""


def _extract_text(path: Path) -> str:
    for encoding in ("utf-8", "latin-1", "cp1252"):
        try:
            return path.read_text(encoding=encoding)
        except UnicodeDecodeError:
            continue
    return ""
